from docx import Document
import pdfplumber
from pptx import Presentation
import os

def extract_text_from_file(file_path, file_type=None):
    if file_type == 'pdf':
        # Extract text from PDF
        return extract_text_from_pdf(file_path)
    elif file_type == 'docx':
        # Extract text from DOCX
        return extract_text_from_docx(file_path)
    elif file_type == 'pptx':
        # Extract text from PPTX
        return extract_text_from_pptx(file_path)
    else:
        # Default or unsupported file type
        raise ValueError("Unsupported file type")

def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path):
    text = ""
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
